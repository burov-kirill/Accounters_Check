import collections
import collections.abc
import datetime

import pandas as pd
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION, XL_TICK_MARK
from pptx.chart.data import CategoryChartData


PLOT_DICT = {'main_table': '',
             'summary_percent': 'Общее выполнение плана по срокам закрытия, %',
             'superior_percent': 'Доля компаний, нарушающих регламентированные',
             'date_table': 'Распределение дат итогового закрытия отчётного месяца',
             'term_data': 'Структурный анализ итогового закрытия отчётного месяца',
             'plan_table': 'Структурный анализ соблюдения Регламентируемых сроков',
             'best_company': 'ТОП-10 «Лидеры закрытия»',
             'worst_company': 'ТОП-10 «Зона роста»',
             'average_table': 'Рейтинг закрытия участков РСБУ на основе балльной системы'}

HUMAN_DATES = {'Январь': 'Января',
               'Февраль': 'Февраля',
               'Март': 'Марта',
               'Апрель': 'Апреля',
               'Май': 'Мая',
               'Июнь': 'Июня',
               'Июль': 'Июля',
               'Август': 'Августа',
               'Сентябрь': 'Сентября',
               'Октябрь': 'Октября',
               'Ноябрь': 'Ноября',
               'Декабрь': 'Декабря'
               }


def set_heading(slide, heading):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(6), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = heading
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = 'FHIDUD+CoFoSans-Medium'
def split_frame(frame, is_dates = False):
    if not is_dates:
        return list(map(lambda x: str(x), list(frame.iloc[:, 0]))),\
            list(map(lambda x: round(x, 2), list(frame.iloc[:, 1])))
    else:
        return list(map(lambda x: str(x), list(frame.iloc[:, 2]))), \
            list(map(lambda x: round(x, 2), list(frame.iloc[:, 1]))), \
                list(map(lambda x: x, list(frame.iloc[:, 0])))
def apply_data_labels(chart):
    plot = chart.plots[0]
    plot.has_data_labels = True
    for series in plot.series:
        values = series.values
        counter = 0
        for point in series.points:
            data_label = point.data_label
            data_label.has_text_frame = True
            data_label.text_frame.text = str(values[counter])
            counter = counter + 1

def set_background(slide):
    left = top = Inches(0)
    background = 'images/background.png'
    pic = slide.shapes.add_picture(background, left, top, width=Inches(16), height=Inches(9))
    slide.shapes._spTree.insert(2, pic._element)

def second_list(root, df, plot_name):
    name = 'Общее выполнение плана по срокам закрытия, %'
    slide = root.slides.add_slide(root.slide_layouts[6])
    set_background(slide)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(6), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = 'Соблюдение Регламентируемых сроков закрытия'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = 'FHIDUD+CoFoSans-Medium'
    # p.line.fill.solid()
    # p.line.fill.fore_color.rgb = RGBColor.from_string('007bfb')
    chart_data = CategoryChartData()

    categories, values = split_frame(df)
    chart_data.categories = categories
    chart_data.add_series(plot_name, values,
                          number_format='0%')
    x, y, cx, cy = Inches(4), Inches(1.5), Inches(8), Inches(6)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = plot_name
    chart.chart_title.text_frame.paragraphs[0].font.fill.solid()
    chart.chart_title.text_frame.paragraphs[0].font.fill.fore_color.rgb = RGBColor.from_string('007bfb')
    chart.has_legend = False
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.has_minor_gridlines = False
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.has_minor_gridlines = False
    chart.category_axis.major_tick_mark = XL_TICK_MARK.NONE
    chart.value_axis.visible = False
    color_list = ["CFCFCF", "007bfb", "ffffff"]
    chart.plots[0].gap_width = 30
    chart.plots[0].has_data_labels = True
    apply_data_labels(chart)

    # Go through every point of the first serie and modify the color
    for idx, point in enumerate(chart.series[0].points):
        col_idx = idx % len(color_list)
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor.from_string(color_list[col_idx])
        point.format.width = Pt(5)
        if col_idx == 2:
            point.format.line.fill.solid()
            point.format.line.fill.fore_color.rgb = RGBColor.from_string(color_list[1])
            point.format.line.width = Pt(2)
def fourth_slide(root, df, plot_name, df2, plot_name2, month, close_date):
    slide = root.slides.add_slide(root.slide_layouts[6])
    set_background(slide)
    set_heading(slide, f'Сводные данные по итоговому закрытию {month}')
    chart_data = CategoryChartData()
    categories, values = split_frame(df2)
    chart_data.categories = categories
    chart_data.add_series(plot_name2, values)
    x, y, cx, cy = Inches(8), Inches(1.5), Inches(8), Inches(6)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data
    ).chart
    # tf = chart.chart_title.text_frame
    # p = tf.add_paragraph()
    # p.font.size = Pt(14)
    # p.font.bold = True
    # p.font.name = 'FHIDUD+CoFoSans-Medium'
    # p.font.color.rgb = RGBColor(255, 255, 255)
    # # p.font.fill.color.rgb = RGBColor(0, 123, 251)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.show_percentage = True
    data_labels.number_format = "0%"
    apply_data_labels(chart)
    color_list = ["00FF00", "ffff00", "FF0000"]
    # Go through every point of the first serie and modify the color
    for idx, point in enumerate(chart.series[0].points):
        col_idx = idx % len(color_list)
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor.from_string(color_list[col_idx])

    chart_data = CategoryChartData()
    categories, values, dates = split_frame(df, is_dates = True)
    chart_data.categories = categories
    chart_data.add_series(plot_name, values)
    x, y, cx, cy = Inches(0.1), Inches(1.5), Inches(8), Inches(6)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = False
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.has_minor_gridlines = False
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.has_minor_gridlines = False

    chart.category_axis.format.line.color.rgb = RGBColor(255, 255, 255)
    chart.value_axis.format.line.color.rgb = RGBColor(255, 255, 255)

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.show_percentage = True
    data_labels.number_format = "0%"
    apply_data_labels(chart)
    color_list = ["00FF00", "ffff00", "FF0000"]
    # Go through every point of the first serie and modify the color
    for idx, (point, dt) in enumerate(zip(chart.series[0].points, dates)):
        if type(dt) == datetime.date:
            diff = close_date - dt
            if diff.days == 0:
                col_idx = 0
            elif diff.days > 0:
                col_idx = 1
            else:
                col_idx = 2
        else:
            col_idx = 2
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor.from_string(color_list[col_idx])
def fifth_slide(root, df, plot_name, month):
    slide = root.slides.add_slide(root.slide_layouts[6])
    set_background(slide)
    set_heading(slide, f'Общие данные по закрытию {HUMAN_DATES[month]}, включая промежуточные участки')
    chart_data = CategoryChartData()
    categories, values = split_frame(df)
    chart_data.categories = categories
    chart_data.add_series(plot_name, values)
    x, y, cx, cy = Inches(8), Inches(1.5), Inches(8), Inches(6)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.show_percentage = True
    data_labels.number_format = "0%"
    apply_data_labels(chart)
    color_list = ["ffff00", "00FF00", "FF0000"]
    # Go through every point of the first serie and modify the color
    for idx, point in enumerate(chart.series[0].points):
        col_idx = idx % len(color_list)
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor.from_string(color_list[col_idx])

def sixth_slide(root, df1, plot_name1, df2, plot_name2):
    slide = root.slides.add_slide(root.slide_layouts[6])
    set_background(slide)
    set_heading(slide, 'Рейтинги категории «труднозакрываемых» организаций группы Самолёт')
    chart_data = CategoryChartData()
    categories, values = split_frame(df1)
    chart_data.categories = categories
    chart_data.add_series(plot_name1, values)
    x, y, cx, cy = Inches(0.1), Inches(1.5), Inches(8), Inches(6)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = False
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.has_minor_gridlines = False
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.has_minor_gridlines = False

    chart.category_axis.format.line.color.rgb = RGBColor(255, 255, 255)
    chart.value_axis.format.line.color.rgb = RGBColor(255, 255, 255)

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.show_percentage = True
    data_labels.number_format = "0%"
    apply_data_labels(chart)
    color_list = ["00FF00"]
    # Go through every point of the first serie and modify the color
    for idx, point in enumerate(chart.series[0].points):
        col_idx = idx % len(color_list)
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor.from_string(color_list[col_idx])

    chart_data = CategoryChartData()
    categories, values = split_frame(df2)
    chart_data.categories = categories
    chart_data.add_series(plot_name2, values)
    x, y, cx, cy = Inches(8), Inches(1.5), Inches(8), Inches(6)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = False
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.has_minor_gridlines = False
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.has_minor_gridlines = False

    chart.category_axis.format.line.color.rgb = RGBColor(255, 255, 255)
    chart.value_axis.format.line.color.rgb = RGBColor(255, 255, 255)

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.show_percentage = True
    data_labels.number_format = "0%"
    apply_data_labels(chart)
    color_list = ["FF0000"]
    # Go through every point of the first serie and modify the color
    for idx, point in enumerate(chart.series[0].points):
        col_idx = idx % len(color_list)
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor.from_string(color_list[col_idx])

def seventh_slide(root, df, plot_name):
    # data = sorted((2.15, 2.13, 2, 1.85, 1.75, 1.54, 1.43, 1.25, 1.13, 1), reverse=True)
    slide = root.slides.add_slide(root.slide_layouts[6])
    set_background(slide)
    set_heading(slide, 'Анализ закрытия участков РСБУ')
    chart_data = CategoryChartData()
    categories, values = split_frame(df)
    chart_data.categories = categories
    chart_data.add_series(plot_name, values)
    x, y, cx, cy = Inches(2.5), Inches(1), Inches(12), Inches(7)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.has_legend = False
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.has_minor_gridlines = False
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.has_minor_gridlines = False

    chart.category_axis.format.line.color.rgb = RGBColor(255, 255, 255)
    chart.value_axis.format.line.color.rgb = RGBColor(255, 255, 255)

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.show_percentage = True
    apply_data_labels(chart)
    color_list = ["ffff00", "00FF00", "FF0000"]
    # Go through every point of the first serie and modify the color
    for point, dt in zip(chart.series[0].points, values):
        if dt>=2:
            color_idx = 0
        elif dt<2 and dt>1.65:
            color_idx = 1
        else:
            color_idx = 2
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor.from_string(color_list[color_idx])

def height_slide(root, month):
    slide = root.slides.add_slide(root.slide_layouts[6])
    set_background(slide)
    set_heading(slide, f'Основные триггеры закрытия {HUMAN_DATES[month]}')
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(6), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = 'Основные триггеры закрытия Июня'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = 'FHIDUD+CoFoSans-Medium'
    p.font.color.rgb = RGBColor(0, 0, 0)
# def edit_tables_dcit(tables, month,  first_table, second_table):
#     temp_frame = pd.DataFrame([[month, tables['summary_percent']]], columns=['month', 'percent'])
#     first_table = pd.concat([first_table, temp_frame])
#     tables['summary_percent'] = first_table
#
#     temp_frame = pd.DataFrame([[month, tables['superior_percent']]], columns=['month', 'percent'])
#     second_table = pd.concat([second_table, temp_frame])
#     tables['superior_percent'] = second_table
#
#     return tables
def create_presentation(path, month, file):
    tables_dict = file.tables
    img_path = 'images/back.png'
    last_img = 'images/last_back.png'
    root = Presentation()
    first_slide_layout = root.slide_layouts[6]
    slide = root.slides.add_slide(first_slide_layout)
    root.slide_height = Inches(9)
    root.slide_width = Inches(16)
    left = top = Inches(0)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(6), Inches(2))
    tf = txBox.text_frame
    pic = slide.shapes.add_picture(img_path,left, top, width=Inches(16), height=Inches(9))
    p = tf.add_paragraph()
    p.text = 'Мониторинг \nзакрытия июня'
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.name = 'FHIDUD+CoFoSans-Medium'
    p.font.color.rgb = RGBColor(255, 255, 255)
    slide.shapes._spTree.insert(2, pic._element)
    closed_date = file.DATE_DICT[file.LAST_COLUMN]

    second_list(root, tables_dict['summary_percent'], PLOT_DICT['summary_percent'])
    second_list(root,tables_dict['superior_percent'], PLOT_DICT['superior_percent'])
    fourth_slide(root,tables_dict['date_table'], PLOT_DICT['date_table'],  tables_dict['term_data'], PLOT_DICT['term_data'], month, closed_date)
    fifth_slide(root, tables_dict['plan_table'], PLOT_DICT['plan_table'], month)
    sixth_slide(root, tables_dict['worst_company'], PLOT_DICT['worst_company'], tables_dict['best_company'], PLOT_DICT['best_company'])
    seventh_slide(root, tables_dict['average_table'], PLOT_DICT['average_table'])
    height_slide(root, month)


    slide = root.slides.add_slide(first_slide_layout)
    left = top = Inches(0)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(2))
    tf = txBox.text_frame
    pic = slide.shapes.add_picture(last_img,left, top, width=Inches(16), height=Inches(9))
    p = tf.add_paragraph()
    p.text = 'СПАСИБО'
    p.font.size = Pt(43)
    p.font.bold = True
    p.font.name = 'FHIDUD+CoFoSans-Medium'
    p.font.color.rgb = RGBColor(255, 255, 255)


    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(6), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = 'По всем вопросам, связанным с мониторингом закрытия\nпериода, можно обратиться к:\n\n' \
             'Любовь Данилина – руководитель кластера Базовый учёт;\n' \
             'Ольга Иванова – методолог РСБУ;\n' \
             'Ольга Клыч  – менеджер внедрения изменений;\n' \
             'Ник Абэ  – координатор, автор отчёта.'
    p.font.size = Pt(18)
    p.font.bold = False
    p.font.name = 'FHIDUD+CoFoSans-Medium'
    p.font.color.rgb = RGBColor(255, 255, 255)
    slide.shapes._spTree.insert(2, pic._element)
    save_path = fr'{path}\Отчет по закрытию.pptx'
    root.save(save_path)

